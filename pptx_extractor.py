from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_slides(pptx_path):
    powerpoint = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(powerpoint.slides, start=1):
        slide_text = []
        slide_images = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
                print(f"Slide {i} text:", shape.text)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                slide_images.append(image.blob)
        slides_data.append({"text": "\n".join(slide_text), "images": slide_images})
    return slides_data

slides_data = extract_slides(r"C:\Users\kyl3g\Downloads\4 Fertilization to Gastrulation FA25.pdf")